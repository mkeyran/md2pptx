import pytest
from click.testing import CliRunner
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE, PP_PARAGRAPH_ALIGNMENT
from pptx.util import Pt

from md2pptx import create_slide, main, parse_content


@pytest.fixture
def runner():
    return CliRunner()


def test_cli(runner, tmp_path):
    input_file = tmp_path / "input.md"
    output_file = tmp_path / "output.pptx"

    input_content = """
# Slide 1
- Bullet 1
- Bullet 2

---

# Slide 2
1. Numbered 1
2. Numbered 2
"""

    input_file.write_text(input_content)

    result = runner.invoke(main, [str(input_file), str(output_file)])
    assert result.exit_code == 0
    assert f"Presentation saved as {output_file}" in result.output

    # Check if the output file was created
    assert output_file.exists()


def test_parse_content():
    content = """
# Header 1
## Header 2
- Bullet 1
- Bullet 2
1. Numbered 1
2. Numbered 2
![Image](https://example.com/image.jpg)
Regular text
- Bullet 3
- Bullet 4
1. Numbered 3
2. Numbered 4
"""
    parsed = parse_content(content)

    assert parsed[0] == ("header", 1, "Header 1")
    assert parsed[1] == ("header", 2, "Header 2")
    assert parsed[2] == ("bullet_list", [(0, "Bullet 1"), (0, "Bullet 2")])
    assert parsed[3] == ("numbered_list", [(0, "Numbered 1"), (0, "Numbered 2")])
    assert parsed[4] == ("image", "https://example.com/image.jpg")
    assert parsed[5] == ("text", "Regular text")
    assert parsed[6] == ("bullet_list", [(0, "Bullet 3"), (0, "Bullet 4")])
    assert parsed[7] == ("numbered_list", [(0, "Numbered 3"), (0, "Numbered 4")])


def test_create_slide():
    prs = Presentation()
    title = "Test Slide"
    content = """
- Bullet 1
  - Nested Bullet
- Bullet 2
1. Numbered 1
2. Numbered 2
   1. Nested Numbered
Regular text
"""
    slide = create_slide(prs, title, content)

    assert slide.shapes.title.text == "Test Slide"
    text_frame = slide.shapes.placeholders[1].text_frame

    assert text_frame.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    paragraphs = text_frame.paragraphs

    assert paragraphs[0].text == "• Bullet 1"
    assert paragraphs[0].level == 0
    assert paragraphs[0].font.size == Pt(18)

    assert paragraphs[1].text == "• Nested Bullet"
    assert paragraphs[1].level == 1
    assert paragraphs[1].font.size == Pt(16)

    assert paragraphs[2].text == "• Bullet 2"
    assert paragraphs[2].level == 0
    assert paragraphs[2].font.size == Pt(18)

    assert paragraphs[3].text == "1. Numbered 1"
    assert paragraphs[3].level == 0
    assert paragraphs[3].font.size == Pt(18)

    assert paragraphs[4].text == "2. Numbered 2"
    assert paragraphs[4].level == 0
    assert paragraphs[4].font.size == Pt(18)

    assert paragraphs[5].text == "1. Nested Numbered"
    assert paragraphs[5].level == 1
    assert paragraphs[5].font.size == Pt(16)

    assert paragraphs[6].text == "Regular text"
    assert paragraphs[6].level == 0
    assert paragraphs[6].alignment == PP_PARAGRAPH_ALIGNMENT.LEFT


# Add more tests as needed for other functions
