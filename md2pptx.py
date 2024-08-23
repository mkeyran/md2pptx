import re
from io import BytesIO

import click
import requests
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.util import Inches, Pt


def add_image_to_slide(
    slide, img_url, left=Inches(1), top=Inches(2.5), width=Inches(6)
):
    try:
        response = requests.get(img_url)
        img_stream = BytesIO(response.content)
        slide.shapes.add_picture(img_stream, left, top, width=width)
    except Exception as e:
        print(f"Error adding image: {e}")


def parse_content(content):
    parsed = []
    current_list = None
    current_list_type = None
    for line in content.split("\n"):
        if line.startswith("!["):
            # Image
            if current_list is not None:
                parsed.append((f"{current_list_type}_list", current_list))
                current_list = None
                current_list_type = None
            img_match = re.search(r"\((.*?)\)", line)
            if img_match:
                parsed.append(("image", img_match.group(1)))
        elif re.match(r"^#{1,6}\s", line):
            # Header
            if current_list is not None:
                parsed.append((f"{current_list_type}_list", current_list))
                current_list = None
                current_list_type = None
            level = len(re.match(r"^#+", line).group(0))
            text = line.lstrip("#").strip()
            parsed.append(("header", level, text.strip()))
        elif re.match(r"^(\s*\d+\.|\s*-|\*)\s", line):
            # List item (numbered or bullet)
            print(line)
            level = len(re.match(r"^\s*", line).group(0)) // 2
            list_type = "numbered" if re.match(r"^\s*\d+\.", line) else "bullet"
            print(list_type)
            text = re.sub(r"^(\s*\d+\.|\s*-|\*)\s", "", line)

            if current_list is None or current_list_type != list_type:
                if current_list is not None:
                    parsed.append((f"{current_list_type}_list", current_list))
                current_list = []
                current_list_type = list_type

            current_list.append((level, text.strip()))
        else:
            # Regular text
            if current_list is not None:
                parsed.append((f"{current_list_type}_list", current_list))
                current_list = None
                current_list_type = None
            if line:
                parsed.append(("text", line.strip()))

    if current_list is not None:
        parsed.append((f"{current_list_type}_list", current_list))

    return parsed


def create_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1]  # Using the bullet slide layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title

    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()  # Clear existing text in the placeholder
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    parsed_content = parse_content(content)

    for item in parsed_content:
        if item[0] == "image":
            add_image_to_slide(slide, item[1])
        elif item[0] == "header":
            level, text = item[1], item[2]
            p = tf.paragraphs[-1]
            p.text = text
            p.level = level - 1  # Adjust level for PowerPoint
            p.font.bold = True
            tf.add_paragraph()
        elif item[0] in ["bullet_list", "numbered_list"]:
            last_levels = {}
            for _, (level, text) in enumerate(item[1], start=1):
                last_levels[level] = last_levels.get(level, 1)
                p = tf.paragraphs[-1]
                p.text = text
                p.level = level
                if item[0] == "numbered_list":
                    p.text = f"{last_levels[level]}. {p.text}"
                    last_levels[level] += 1
                else:
                    p.text = f"• {p.text}"

                # Adjust font size based on level
                if level == 0:
                    p.font.size = Pt(18)
                elif level == 1:
                    p.font.size = Pt(16)
                else:
                    p.font.size = Pt(14)
                tf.add_paragraph()
        elif item[0] == "text":
            p = tf.paragraphs[-1]
            p.text = item[1]
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            tf.add_paragraph()

    return slide


def markdown_to_pptx(prs, markdown_content, output_file):
    slides = re.split(r"\n---\n", markdown_content)

    for slide_content in slides:
        parts = slide_content.strip().split("\n", 1)
        if len(parts) > 1:
            title, content = parts
        else:
            title, content = parts[0], ""

        title = title.lstrip("#").strip()
        create_slide(prs, title, content)

    prs.save(output_file)


@click.command()
@click.argument("input_file", type=click.Path(exists=True))
@click.argument("output_file", type=click.Path())
@click.option("--width", default=16.0, help="Slide width in inches", type=float)
@click.option("--height", default=9.0, help="Slide height in inches", type=float)
def main(input_file, output_file, width, height):
    """Convert Markdown to PowerPoint presentation."""
    with open(input_file, "r", encoding="utf-8") as f:
        markdown_content = f.read()

    prs = Presentation()
    prs.slide_width = Inches(width)
    prs.slide_height = Inches(height)

    slides = re.split(r"\n---\n", markdown_content)

    for slide_content in slides:
        parts = slide_content.strip().split("\n", 1)
        if len(parts) > 1:
            title, content = parts
        else:
            title, content = parts[0], ""

        title = title.lstrip("#").strip()
        create_slide(prs, title, content)

    prs.save(output_file)
    click.echo(f"Presentation saved as {output_file}")


if __name__ == "__main__":
    main()
